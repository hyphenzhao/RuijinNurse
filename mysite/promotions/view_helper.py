from pathlib import Path
from datetime import datetime

from django import forms
from django.shortcuts import render
from django.conf import settings
import subprocess, shutil, re, os, json, time, requests
from django.http import JsonResponse, StreamingHttpResponse
from django.views.decorators.csrf import csrf_exempt, ensure_csrf_cookie, csrf_protect
from django.utils.encoding import smart_str

# Optional imports — only needed by specific functions (not used by the
# shared SSE generator or the mobile API).
try:
    from openai import OpenAI
except ImportError:  # pragma: no cover
    OpenAI = None

try:
    import qianfan
except ImportError:  # pragma: no cover
    qianfan = None


def _safe_model_key(model_key: str) -> str:
    """Sanitize model key for filesystem use."""
    return re.sub(r'[^A-Za-z0-9._-]+', '_', model_key or 'default')


def _resolve_session_identifier(request_or_user_id):
    """
    Return a stable session identifier for session-context persistence.

    If given a Django request with an authenticated user, returns 'user_{id}'.
    If given a raw user ID (int/str), returns 'user_{id}'.
    Otherwise falls back to the Django session key.
    """
    # If passed a request object with an authenticated user
    if hasattr(request_or_user_id, 'user') and request_or_user_id.user and request_or_user_id.user.is_authenticated:
        return f'user_{request_or_user_id.user.id}'

    # If passed a request object without auth, fall back to session key
    if hasattr(request_or_user_id, 'session'):
        session_key = request_or_user_id.session.session_key
        if not session_key:
            request_or_user_id.session.save()
            session_key = request_or_user_id.session.session_key
        return session_key

    # If passed a raw int or digit string (user ID)
    uid = str(request_or_user_id)
    return f'user_{uid}'


def generate_sse_stream(
    *,
    chat_target: dict,
    default_role_text: str,
    default_knowledge_text: str,
    promotion: str,
    model_key: str,
    session_path: Path,
    chat_history_callback=None,
):
    """
    Shared SSE stream generator used by both the web view and mobile API.

    Yields SSE event strings.  On completion, saves the session context
    to *session_path* and calls *chat_history_callback(model, promotion, full_text)* if given.
    """
    helper_sse = _helper_sse

    full_text = ''
    think_text = ''
    think_state = {'in_think': False}

    try:
        role_text = chat_target['role_text'] if chat_target['role_text'] is not None else default_role_text
        knowledge_text = chat_target['knowledge_text'] if chat_target['knowledge_text'] is not None else default_knowledge_text
        system_context = _build_system_context(role_text, knowledge_text)

        session_data = _load_session_context_from_path(session_path, system_context, model_key)
        session_data['system_context'] = system_context

        with _stream_ollama_response_by_target(
            chat_target['base_url'],
            chat_target['model_name'],
            _build_messages(system_context, _history_messages_for_ollama(session_data), promotion),
        ) as response:
            response.raise_for_status()
            for line in response.iter_lines(decode_unicode=True):
                if not line:
                    continue
                try:
                    data = json.loads(line)
                except json.JSONDecodeError:
                    continue

                message = data.get('message', {}) or {}

                # Newer Ollama thinking API: reasoning is streamed via message.thinking
                direct_thinking_delta = message.get('thinking', '')
                if direct_thinking_delta:
                    think_text += direct_thinking_delta
                    yield helper_sse('thinking', {'content': direct_thinking_delta})

                # Final answer text
                delta = message.get('content', '')

                # Backward compatibility: some models emit <think>...</think> tags
                if delta:
                    thinking_delta, answer_delta = _split_thinking_delta(delta, think_state)
                else:
                    thinking_delta, answer_delta = '', ''

                if thinking_delta:
                    think_text += thinking_delta
                    yield helper_sse('thinking', {'content': thinking_delta})

                if answer_delta:
                    full_text += answer_delta
                    yield helper_sse('delta', {'content': answer_delta})
                elif delta and not direct_thinking_delta:
                    full_text += delta
                    yield helper_sse('delta', {'content': delta})

                if data.get('done', True):
                    break

        # Save session context
        _append_session_message(session_data, 'user', promotion)
        _append_session_message(session_data, 'assistant', full_text)
        _save_session_context(session_path, session_data)

    except Exception as e:
        base_url = chat_target.get('base_url', 'unknown')
        err = f'（服务端错误）无法连接 Ollama ({base_url}): {type(e).__name__}: {e}'
        yield helper_sse('delta', {'content': err})
        full_text = (full_text or '') + '\n' + err

    if chat_history_callback:
        try:
            chat_history_callback(model_key, promotion, full_text)
        except Exception:
            pass

    yield helper_sse('done', {'final': full_text, 'thinking': think_text, 'audio_url': None})


# ---- Internal helpers used by generate_sse_stream (namespace-imported above) ----

def _helper_sse(event: str, data: dict) -> str:
    return f"event: {event}\ndata: {json.dumps(data, ensure_ascii=False)}\n\n"


def _build_system_context(role_text: str, knowledge_text: str) -> str:
    role_text = (role_text or '').strip()
    knowledge_text = (knowledge_text or '').strip()
    parts = []
    if role_text:
        parts.append(f'[角色设定]\n{role_text}')
    if knowledge_text:
        parts.append(f'[参考知识]\n{knowledge_text}')
    return '\n\n'.join(parts).strip()


def _build_messages(system_context: str, history_messages, promotion: str):
    messages = []
    if system_context:
        messages.append({'role': 'system', 'content': system_context})
    if history_messages:
        messages.extend(history_messages)
    messages.append({'role': 'user', 'content': promotion})
    return messages


def _build_ollama_base_url(host: str, port) -> str:
    return f'http://{str(host).strip()}:{str(port).strip()}'


def _stream_ollama_response_by_target(base_url: str, model_name: str, messages):
    import requests as _requests
    payload = {
        'model': model_name,
        'messages': messages,
        'think': True,
        'stream': True,
    }
    return _requests.post(f'{base_url}/api/chat', json=payload, stream=True)


def _split_thinking_delta(delta: str, state: dict):
    """Split <think>...</think> from answer text in one chunk."""
    if not delta:
        return '', ''

    thinking_parts = []
    answer_parts = []
    buf = delta

    while buf:
        if state['in_think']:
            end = buf.find('</think>')
            if end == -1:
                thinking_parts.append(buf)
                buf = ''
            else:
                thinking_parts.append(buf[:end])
                buf = buf[end + len('</think>'):]
                state['in_think'] = False
            continue

        start = buf.find('<think>')
        if start == -1:
            answer_parts.append(buf)
            buf = ''
        else:
            if start > 0:
                answer_parts.append(buf[:start])
            buf = buf[start + len('<think>'):]
            state['in_think'] = True

    return ''.join(thinking_parts), ''.join(answer_parts)


def _load_session_context_from_path(path: Path, system_context: str, model_key: str):
    if path.exists():
        try:
            data = json.loads(path.read_text(encoding='utf-8'))
            if isinstance(data, dict):
                data.setdefault('version', 1)
                data.setdefault('model_key', model_key)
                data.setdefault('system_context', system_context)
                data.setdefault('messages', [])
                return data
        except Exception:
            pass
    return {
        'version': 1,
        'model_key': model_key,
        'system_context': system_context,
        'messages': [],
        'updated_at': None,
    }


def _save_session_context(path: Path, data: dict):
    data['updated_at'] = datetime.now().isoformat()
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding='utf-8')


def _append_session_message(data: dict, role: str, content: str):
    content = (content or '').strip()
    if not content:
        return
    data.setdefault('messages', [])
    data['messages'].append({
        'role': role,
        'content': content,
        'timestamp': datetime.now().isoformat(),
    })


def _history_messages_for_ollama(data: dict):
    out = []
    for item in data.get('messages', []):
        role = item.get('role')
        content = item.get('content')
        if role in {'user', 'assistant'} and content:
            out.append({'role': role, 'content': content})
    return out

command = "./tts_offline_sample"
working_directory = "/Volumes/Workspace/RuijinNurse/mysite/promotions/static/promotions/"
STATIC_WAV_PATH = '/Volumes/Workspace/RuijinNurse/mysite/promotions/static/promotions/tts_sample.wav'
# STATIC_PATH     = '/Volumes/HDWD16TB/Workspace/RuijinNurse/mysite/promotions/static/promotions/'


def load_role_and_knowledge(base_path):
    role_file_path = os.path.join(base_path, 'role.txt')
    knowledge_file_path = os.path.join(base_path, 'knowledge.txt')

    # Initialize variables in case files don't exist
    role_text = ''
    knowledge_text = ''
    if os.path.exists(role_file_path):
        with open(role_file_path, 'r', encoding='utf-8') as f:
            role_text = f.read()
    if os.path.exists(knowledge_file_path):
        with open(knowledge_file_path, 'r', encoding='utf-8') as f:
            knowledge_text = f.read()
    return role_text, knowledge_text

# def run_tts_program():
#     try:
#         result = subprocess.run(
#             command, 
#             cwd=working_directory,  # Set the working directory
#             check=True,  # Raises an exception if the command fails
#             stdout=subprocess.PIPE,  # Capture standard output
#             stderr=subprocess.PIPE   # Capture standard error
#         )

#         # Print the output of the command
#         print("Output:", result.stdout.decode())
#         print("Error:", result.stderr.decode())

#     except subprocess.CalledProcessError as e:
#         print("An error occurred while running the command.")
#         print("Error code:", e.returncode)
#         print("Error message:", e.stderr.decode())
def split_chinese_sentences(text):
    # The regex pattern splits on Chinese punctuation that is used for ending sentences.
    # It uses a capturing group so that the punctuation is preserved.
    parts = re.split(r'([。！？：；])', text)
    
    sentences = []
    sentence = ""
    # Reassemble the text parts and punctuation into full sentences.
    for part in parts:
        # Append the part (could be text or punctuation) to the current sentence.
        sentence += part.strip()
        # If the part is a Chinese punctuation mark that ends a sentence,
        # we consider the sentence complete.
        if re.fullmatch(r'[。！？]', part):
            sentences.append(sentence)
            sentence = ""
    # If there's any remaining text without a trailing punctuation, add it as well.
    if sentence:
        sentences.append(sentence)
    return sentences

def run_tts_program(generated_audio, full_content):
    # Define the OpenTTS URL; note that voice can be passed either in the URL or in the payload.
    opentts_url = "http://localhost:5500/api/tts?voice=coqui-tts:zh_baker"
    # opentts_url = "http://localhost:5500/api/tts?voice=espeak:zh"

    # Build your payload with the Chinese text.
    payload = {
        "text": split_chinese_sentences(full_content)  # e.g., "你好"
    }
    
    # Convert payload to a JSON string using ensure_ascii=False so that non-ASCII characters remain intact.
    payload_json = json.dumps(full_content, ensure_ascii=False)
    
    # Set the header to indicate UTF-8 encoding.
    headers = {"Content-Type": "application/json; charset=utf-8"}

    # Send the POST request using data=payload_json to use your custom JSON encoding.
    response = requests.post(opentts_url, data=full_content.replace("；", "。").replace("：", "。"), headers=headers)
    
    # Debug: print the response content for verification
    # print("Response content:", response.content)
    
    # Check the response status code.
    if response.status_code == 200:
        # Optionally check that the returned content type is audio/wav.
        content_type = response.headers.get("Content-Type", "")
        if "audio/wav" in content_type:
            with open(generated_audio, "wb") as audio_file:
                audio_file.write(response.content)
            print(f"Audio saved successfully to: {generated_audio}")
        else:
            print("Unexpected Content-Type:", content_type)
    else:
        print(f"Failed to generate audio. Status code: {response.status_code}\nResponse: {response.text}")

def convert_wav_to_mp3(input_path, output_path):
    # Run ffmpeg to convert wav to mp3
    command = ["ffmpeg", "-i", input_path, "-q:a", "2", output_path]
    try:
        subprocess.run(command, check=True)
        print("Conversion to MP3 successful.")
    except subprocess.CalledProcessError as e:
        print("An error occurred during conversion:", e)

def helper_sse(event: str, data: dict) -> str:
    """Pack an SSE-like event line."""
    return f"event: {event}\ndata: {json.dumps(data, ensure_ascii=False)}\n\n"

def helper_chunk_iter(text: str, size: int = 24):
    """Yield small chunks to simulate/standardize streaming for non-streaming APIs."""
    for i in range(0, len(text), size):
        yield text[i:i+size]

def helper_filter_think_stream(delta: str, state: dict) -> str:
    """
    Remove any <think>...</think> segments on-the-fly, maintaining state across chunks.
    state = {'in_think': bool}
    """
    if not delta:
        return ""

    out = []
    buf = delta

    while buf:
        if state['in_think']:
            end = buf.find("</think>")
            if end == -1:
                # still inside think, drop everything
                return ""
            # close tag found; drop it and switch off
            buf = buf[end+len("</think>"):]
            state['in_think'] = False
            continue

        # not currently in think; look for start tag
        start = buf.find("<think>")
        if start == -1:
            # no think tag, all is valid content
            out.append(buf)
            break
        # append content before <think>
        if start > 0:
            out.append(buf[:start])
        # enter think
        buf = buf[start+len("<think>"):]
        state['in_think'] = True
        # loop to try to find </think> in remaining buf

    return "".join(out)

def extract_text_from_upload(uploaded_file):
    """
    根据文件扩展名，从 Word / PPTX / PDF 中抽取纯文本。
    仅支持: .docx, .pptx, .pdf
    """
    import os
    ext = os.path.splitext(uploaded_file.name)[1].lower()

    if ext == '.docx':
        from docx import Document
        doc = Document(uploaded_file)
        return "\n".join(p.text for p in doc.paragraphs)

    elif ext == '.pptx':
        from pptx import Presentation
        prs = Presentation(uploaded_file)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)

    elif ext == '.pdf':
        from PyPDF2 import PdfReader
        reader = PdfReader(uploaded_file)
        texts = []
        for page in reader.pages:
            txt = page.extract_text() or ""
            texts.append(txt)
        return "\n".join(texts)

    elif ext == '.txt':
        return uploaded_file.read().decode('utf-8', errors='replace')

    else:
        raise ValueError(f"不支持的文件类型: {ext}")


def extract_files_from_upload(uploaded_files):
    """Extract text from multiple uploaded files and/or zip archives.

    Returns a list of dicts: [{'filename': str, 'text': str, 'error': str|None}]
    """
    import zipfile, io, os, tempfile

    results = []

    for uploaded_file in uploaded_files:
        fname = uploaded_file.name.lower()
        try:
            if fname.endswith('.zip'):
                # Extract zip and process each file inside
                with zipfile.ZipFile(uploaded_file) as zf:
                    for member in zf.namelist():
                        # Skip directories and hidden files
                        if member.endswith('/') or os.path.basename(member).startswith('.'):
                            continue
                        ext = os.path.splitext(member)[1].lower()
                        if ext not in ('.pdf', '.docx', '.pptx', '.txt', '.doc', '.ppt'):
                            continue
                        try:
                            inner_bytes = zf.read(member)
                            inner_text = _extract_bytes(inner_bytes, os.path.basename(member))
                            results.append({'filename': member, 'text': inner_text, 'error': None})
                        except Exception as e:
                            results.append({'filename': member, 'text': '', 'error': str(e)})
            else:
                text = extract_text_from_upload(uploaded_file)
                results.append({'filename': uploaded_file.name, 'text': text, 'error': None})
        except Exception as e:
            results.append({'filename': uploaded_file.name, 'text': '', 'error': str(e)})

    return results


def _extract_bytes(data: bytes, filename: str) -> str:
    """Extract text from in-memory bytes given a filename (used for zip members)."""
    import io, os
    ext = os.path.splitext(filename)[1].lower()

    if ext == '.txt':
        return data.decode('utf-8', errors='replace')
    elif ext == '.pdf':
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(data))
        return "\n".join(p.extract_text() or "" for p in reader.pages)
    elif ext == '.docx':
        from docx import Document
        doc = Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs)
    elif ext == '.pptx':
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    else:
        return data.decode('utf-8', errors='replace')


# ══════════════════════════════════════════════════════════════
#  Knowledge Graph RAG — intent routing + retrieval
# ══════════════════════════════════════════════════════════════

def route_and_retrieve(agent, user_query):
    """Two-stage routing for knowledge-graph agents.

    Stage 1: keyword match against ReferenceCards.
    Stage 2: semantic search against KnowledgeEntries.

    Returns a dict:
      {'route': 'reference_card', 'card': <ReferenceCard>, 'context': <str>}
      {'route': 'rag', 'context': <str>}
      {'route': 'none', 'context': ''}
    """
    if not agent or agent.agent_type != 'kg':
        return {'route': 'none', 'context': ''}

    kgs = agent.knowledge_graphs.filter(is_active=True)
    if not kgs.exists():
        return {'route': 'none', 'context': ''}

    # Stage 1: try ReferenceCard exact/keyword match across all bound KGs
    best_card = None
    best_score = 0
    for kg in kgs:
        card, score = match_reference_card(kg, user_query)
        if card and score > best_score:
            best_card = card
            best_score = score

    if best_card and best_score > 0.5:
        return {
            'route': 'reference_card',
            'card': best_card,
            'context': format_reference_card(best_card),
        }

    # Stage 2: full-text search on KnowledgeEntries
    knowledge_chunks = []
    for kg in kgs:
        chunks = semantic_search_knowledge(kg, user_query, top_k=3)
        knowledge_chunks.extend(chunks)

    # Also search literature if the query looks academic
    if _is_research_query(user_query):
        lit_chunks = []
        for kg in kgs:
            chunks = semantic_search_literature(kg, user_query, top_k=2)
            lit_chunks.extend(chunks)
        knowledge_chunks.extend(lit_chunks)

    if knowledge_chunks:
        context = _build_rag_context(knowledge_chunks)
        return {'route': 'rag', 'context': context}

    return {'route': 'none', 'context': ''}


def match_reference_card(knowledge_graph, user_query):
    """Keyword-based matching for ReferenceCards.

    Returns (card, confidence) where confidence is 0.0-1.0.
    Higher confidence = more keyword overlap + shorter card = better match.
    """
    from .models import ReferenceCard

    query_lower = user_query.lower()
    query_chars = set(query_lower)

    best_card = None
    best_score = 0.0

    for card in knowledge_graph.reference_cards.filter(is_active=True):
        if not card.trigger_keywords:
            continue

        keywords = [kw.strip().lower() for kw in card.trigger_keywords.split('\n') if kw.strip()]
        if not keywords:
            continue

        match_count = 0
        for kw in keywords:
            if kw in query_lower:
                match_count += 1
            else:
                # Partial character overlap for Chinese
                kw_chars = set(kw)
                if len(kw_chars & query_chars) >= max(2, len(kw_chars) * 0.6):
                    match_count += 0.5

        if match_count == 0:
            continue

        # Score: matches / sqrt(num_keywords) — rewards high match ratio
        score = match_count / max(1, len(keywords) ** 0.5)

        if score > best_score:
            best_score = score
            best_card = card

    return best_card, best_score


def semantic_search_knowledge(knowledge_graph, user_query, top_k=3):
    """Full-text search on KnowledgeEntries.

    Uses simple word-overlap scoring (suitable for Chinese text).
    """
    from .models import KnowledgeEntry

    query_lower = user_query.lower()
    query_words = set(query_lower)

    entries = knowledge_graph.knowledge_entries.filter(is_active=True)
    if not entries.exists():
        return []

    scored = []
    for entry in entries:
        content_lower = entry.content.lower()
        # Word overlap score
        score = len(query_words & set(content_lower))
        # Bonus for title match
        if entry.title.lower() in query_lower or any(w in entry.title.lower() for w in query_words if len(w) > 1):
            score += 3
        # Bonus for tag match
        if entry.tags:
            tag_words = set(entry.tags.lower().replace(',', ' ').split())
            score += len(query_words & tag_words) * 2
        if score > 0:
            scored.append((score, entry))

    scored.sort(key=lambda x: x[0], reverse=True)
    return [{'title': entry.title, 'content': entry.content, 'source': entry.source}
            for _, entry in scored[:top_k]]


def semantic_search_literature(knowledge_graph, user_query, top_k=2):
    """Full-text search on LiteratureEntry abstracts."""
    from .models import LiteratureEntry

    query_lower = user_query.lower()
    query_words = set(query_lower)

    entries = knowledge_graph.literature.filter(is_active=True)
    if not entries.exists():
        return []

    scored = []
    for entry in entries:
        abstract_lower = entry.abstract.lower()
        score = len(query_words & set(abstract_lower))
        title_lower = entry.title.lower()
        score += len(query_words & set(title_lower)) * 2
        if score > 0:
            scored.append((score, entry))

    scored.sort(key=lambda x: x[0], reverse=True)
    return [{'title': entry.title, 'content': entry.abstract,
             'source': f'{entry.journal}, {entry.year}' if entry.year else entry.journal,
             'evidence_level': entry.evidence_level}
            for _, entry in scored[:top_k]]


def format_reference_card(card):
    """Format a ReferenceCard as a structured context string for the LLM."""
    lines = [f'【宣教指南】{card.title}']
    if card.media_type != 'none' and card.media_url:
        video_line = f'[{card.title}]({card.media_url})'
        lines.append(f'📹 宣教视频（请在回复中以此 Markdown 链接格式提供）：{video_line}')
        if card.media_timespan:
            lines.append(f'   视频时间段：{card.media_timespan}')
    if card.key_points:
        lines.append('📋 要点摘要（请在回复中列出）：')
        for kp in card.key_points:
            label = kp.get('label', '')
            text = kp.get('text', '')
            if label:
                lines.append(f'  - {label}：{text}')
            else:
                lines.append(f'  - {text}')
    # Add explicit instruction for the LLM
    if card.media_type != 'none' and card.media_url:
        lines.append('\n⚠️ 重要：请在回复中务必包含上述视频链接（使用 Markdown 链接格式 [标题](URL)），让用户可以直接点击观看。同时列出要点摘要。')
    else:
        lines.append('\n请在回复中列出上述要点摘要，帮助用户快速理解。')
    return '\n'.join(lines)


def _build_rag_context(chunks):
    """Build a combined context string from retrieval results."""
    parts = []
    for i, chunk in enumerate(chunks, 1):
        src = chunk.get('source', '')
        title = chunk.get('title', '')
        content = chunk.get('content', '')
        evidence = chunk.get('evidence_level', '')
        header = f'【参考资料 {i}】{title}'
        if src:
            header += f'（来源：{src}）'
        if evidence:
            header += f' [证据等级：{evidence}]'
        parts.append(f'{header}\n{content}')
    return '\n\n'.join(parts)


def _is_research_query(user_query):
    """Heuristic: does the query look like an academic/research question?"""
    research_keywords = [
        '研究', '文献', '证据', '论文', '临床试验', 'meta', '综述',
        '最新', '进展', '指南', '共识', 'RCT', '系统评价',
    ]
    query_lower = user_query.lower()
    return any(kw in query_lower for kw in research_keywords)


# ══════════════════════════════════════════════════════════════
#  AI Extraction — call Ollama to parse unstructured text
# ══════════════════════════════════════════════════════════════

_EXTRACTION_PROMPTS = {
    'reference_card': """你是一个医疗知识提取助手。请从以下文本中提取"流程指南卡"的结构化信息。
严格返回 JSON 格式（不要 markdown 代码块），字段如下：
{
  "title": "指南标题",
  "trigger_keywords": "关键词1\\n关键词2\\n关键词3",
  "category": "admission/discharge/medication/rehab/diet/safety/followup/general 之一",
  "media_type": "video/image/pdf/none 之一",
  "media_url": "视频或图片链接（如有）",
  "media_timespan": "时间区间如 00:09-01:39（如有）",
  "key_points": [{"label": "要点标题", "text": "要点内容"}]
}""",

    'knowledge_entry': """你是一个医疗知识提取助手。请从以下文本中提取"专业知识条目"的结构化信息。
严格返回 JSON 格式（不要 markdown 代码块），字段如下：
{
  "title": "知识条目标题（简洁概括）",
  "content": "知识内容（保留原文关键信息，适当精炼）",
  "source": "知识来源（如有）",
  "tags": "标签1,标签2,标签3"
}""",

    'literature': """你是一个学术文献分析助手。请从以下文本中提取"学术文献"的结构化信息。
严格返回 JSON 格式（不要 markdown 代码块），字段如下：
{
  "title": "文献标题",
  "authors": "作者（逗号分隔）",
  "year": 年份数字（如有）,
  "journal": "期刊名称（如有）",
  "doi": "DOI（如有）",
  "abstract": "摘要内容",
  "evidence_level": "A/B/C/D 或空字符串"
}""",
}


_BATCH_REFERENCE_CARD_PROMPT = """你是一个医疗知识提取助手。请从以下文本中提取所有"流程指南卡"。
文本可能包含多个独立主题的指南（如住院手续、病区环境、出院流程等），请为每个主题生成一张指南卡。
严格返回 JSON 数组格式（不要 markdown 代码块），每项字段如下：
[{
  "title": "指南标题",
  "trigger_keywords": "关键词1\\n关键词2\\n关键词3",
  "category": "admission/discharge/medication/rehab/diet/safety/followup/general 之一",
  "media_type": "video/image/pdf/none 之一",
  "media_url": "视频或图片链接（如有）",
  "media_timespan": "时间区间如 00:09-01:39（如有）",
  "key_points": [{"label": "要点标题", "text": "要点内容"}]
}]"""


def ai_extract_knowledge_batch(knowledge_graph, raw_text):
    """Batch extract multiple ReferenceCards from raw text."""
    if not knowledge_graph.extraction_model:
        return False, '该知识图谱未配置 AI 提取模型。'

    base_url = f'http://{knowledge_graph.ollama_host}:{knowledge_graph.ollama_port}'
    prompt = f'{_BATCH_REFERENCE_CARD_PROMPT}\n\n---\n待提取文本：\n{raw_text[:10000]}'

    try:
        resp = requests.post(
            f'{base_url}/api/chat',
            json={
                'model': knowledge_graph.extraction_model,
                'messages': [{'role': 'user', 'content': prompt}],
                'stream': False,
            },
            timeout=180,
        )
        resp.raise_for_status()
        data = resp.json()
        content = data.get('message', {}).get('content', '').strip()

        if content.startswith('```'):
            content = re.sub(r'^```(?:json)?\s*', '', content)
            content = re.sub(r'\s*```$', '', content)

        try:
            result = json.loads(content)
        except json.JSONDecodeError:
            match = re.search(r'\[[\s\S]*\]', content)
            if match:
                result = json.loads(match.group(0))
            else:
                return False, f'AI 返回无法解析为 JSON: {content[:200]}...'

        if not isinstance(result, list):
            result = [result]
        return True, result

    except requests.RequestException as e:
        return False, f'连接 Ollama 失败: {type(e).__name__}: {e}'
    except Exception as e:
        return False, f'AI 批量提取失败: {type(e).__name__}: {e}'


def ai_extract_knowledge(knowledge_graph, knowledge_type, raw_text):
    """Call Ollama extraction model to parse unstructured text into structured JSON.

    Args:
        knowledge_graph: KnowledgeGraph instance
        knowledge_type: 'reference_card' | 'knowledge_entry' | 'literature'
        raw_text: raw text to extract from

    Returns:
        (success: bool, data: dict | error_message: str)
    """
    if not knowledge_graph.extraction_model:
        return False, '该知识图谱未配置 AI 提取模型。'

    prompt_template = _EXTRACTION_PROMPTS.get(knowledge_type)
    if not prompt_template:
        return False, f'不支持的知识类型: {knowledge_type}'

    base_url = f'http://{knowledge_graph.ollama_host}:{knowledge_graph.ollama_port}'
    prompt = f'{prompt_template}\n\n---\n待提取文本：\n{raw_text[:8000]}'

    try:
        resp = requests.post(
            f'{base_url}/api/chat',
            json={
                'model': knowledge_graph.extraction_model,
                'messages': [{'role': 'user', 'content': prompt}],
                'stream': False,
            },
            timeout=120,
        )
        resp.raise_for_status()
        data = resp.json()
        content = data.get('message', {}).get('content', '')

        # Try to parse JSON from the response
        content = content.strip()
        # Remove markdown code fences if present
        if content.startswith('```'):
            content = re.sub(r'^```(?:json)?\s*', '', content)
            content = re.sub(r'\s*```$', '', content)

        try:
            result = json.loads(content)
        except json.JSONDecodeError:
            # Try to extract JSON block
            match = re.search(r'\{[\s\S]*\}', content)
            if match:
                result = json.loads(match.group(0))
            else:
                return False, f'AI 返回的内容无法解析为 JSON: {content[:200]}...'

        return True, result

    except requests.RequestException as e:
        return False, f'连接 Ollama 失败: {type(e).__name__}: {e}'
    except Exception as e:
        return False, f'AI 提取失败: {type(e).__name__}: {e}'
